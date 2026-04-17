"""Build VibeX hackathon pitch deck (.pptx).

Style pulled from DESIGN.md: dark panels, neon accents, pixel/retro feel.
Content pulled from README.md.
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
SCREENS = ROOT / "docs" / "screenshots"
OUT = Path("/Users/alexji/Desktop/VibeX-Hackathon-Pitch.pptx")

# --- Palette (from DESIGN.md) -------------------------------------------------
BG_DEEP = RGBColor(0x0D, 0x0D, 0x0D)
BG_PANEL = RGBColor(0x11, 0x11, 0x14)
BG_CARD = RGBColor(0x16, 0x16, 0x19)
BORDER = RGBColor(0x2A, 0x2A, 0x30)
BORDER_BOLT = RGBColor(0x3A, 0x3A, 0x42)

TEXT = RGBColor(0xE8, 0xE8, 0xEC)
MUTED = RGBColor(0x88, 0x88, 0xA0)
DIM = RGBColor(0x55, 0x55, 0x55)

NEON_GREEN = RGBColor(0x39, 0xFF, 0x14)
NEON_PURPLE = RGBColor(0x9D, 0x00, 0xFF)
NEON_ORANGE = RGBColor(0xFF, 0x45, 0x00)
NEON_YELLOW = RGBColor(0xFA, 0xCC, 0x15)
NEON_CYAN = RGBColor(0x06, 0xB6, 0xD4)
NEON_PINK = RGBColor(0xEC, 0x48, 0x99)

# Fonts — fall back to system fonts that feel pixel/retro.
PIXEL_FONT = "Courier New"  # heading / pixel vibe
BODY_FONT = "Consolas"       # body / retro terminal

# --- Slide geometry (16:9) ----------------------------------------------------
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


# --- Helpers ------------------------------------------------------------------
def fill_solid(shape, rgb):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb


def no_line(shape):
    shape.line.fill.background()


def line_solid(shape, rgb, width_pt=2):
    shape.line.color.rgb = rgb
    shape.line.width = Pt(width_pt)


def add_rect(slide, left, top, width, height, fill=None, line=None, line_w=2):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    if fill is None:
        shape.fill.background()
    else:
        fill_solid(shape, fill)
    if line is None:
        no_line(shape)
    else:
        line_solid(shape, line, line_w)
    shape.shadow.inherit = False
    return shape


def add_text(
    slide,
    left,
    top,
    width,
    height,
    text,
    *,
    font=BODY_FONT,
    size=14,
    color=TEXT,
    bold=False,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.bold = bold
    return box


def slide_background(slide, color=BG_DEEP):
    bg = add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=color)
    # Push to back
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)
    return bg


def chrome(slide, eyebrow="VIBEX://PITCH_V1", page=None, total=14):
    # Top-left pixel eyebrow
    add_text(
        slide, Inches(0.5), Inches(0.35), Inches(6), Inches(0.3),
        eyebrow, font=PIXEL_FONT, size=10, color=NEON_GREEN, bold=True,
    )
    # Thin neon bar under eyebrow
    add_rect(slide, Inches(0.5), Inches(0.7), Inches(1.2), Emu(28575),
             fill=NEON_PURPLE)
    # Bottom footer
    add_text(
        slide, Inches(0.5), Inches(7.05), Inches(8), Inches(0.3),
        "VIBEXFORGE.COM", font=PIXEL_FONT, size=9, color=MUTED, bold=True,
    )
    if page is not None:
        add_text(
            slide, Inches(12.0), Inches(7.05), Inches(1.3), Inches(0.3),
            f"{page:02d} / {total:02d}",
            font=PIXEL_FONT, size=9, color=NEON_PURPLE, bold=True,
            align=PP_ALIGN.RIGHT,
        )


def l_corners(slide):
    """Orange L-corner brackets (DESIGN.md §Decorative)."""
    size = Inches(0.22)
    th = Emu(38100)  # ~3pt
    # top-left
    add_rect(slide, Inches(0.2), Inches(0.2), size, th, fill=NEON_ORANGE)
    add_rect(slide, Inches(0.2), Inches(0.2), th, size, fill=NEON_ORANGE)
    # top-right
    add_rect(slide, SLIDE_W - Inches(0.2) - size, Inches(0.2), size, th, fill=NEON_ORANGE)
    add_rect(slide, SLIDE_W - Inches(0.2) - th, Inches(0.2), th, size, fill=NEON_ORANGE)
    # bottom-left
    add_rect(slide, Inches(0.2), SLIDE_H - Inches(0.2) - th, size, th, fill=NEON_ORANGE)
    add_rect(slide, Inches(0.2), SLIDE_H - Inches(0.2) - size, th, size, fill=NEON_ORANGE)
    # bottom-right
    add_rect(slide, SLIDE_W - Inches(0.2) - size, SLIDE_H - Inches(0.2) - th, size, th, fill=NEON_ORANGE)
    add_rect(slide, SLIDE_W - Inches(0.2) - th, SLIDE_H - Inches(0.2) - size, th, size, fill=NEON_ORANGE)


def section_title(slide, eyebrow, title, title_color=TEXT, accent=NEON_PURPLE):
    add_text(
        slide, Inches(0.5), Inches(1.1), Inches(10), Inches(0.35),
        eyebrow, font=PIXEL_FONT, size=11, color=accent, bold=True,
    )
    add_text(
        slide, Inches(0.5), Inches(1.45), Inches(12.3), Inches(0.9),
        title, font=PIXEL_FONT, size=30, color=title_color, bold=True,
    )
    # accent underline
    add_rect(slide, Inches(0.5), Inches(2.35), Inches(0.8), Emu(38100), fill=accent)


def glass_card(slide, left, top, width, height, accent=NEON_PURPLE):
    # 4px top accent bar
    add_rect(slide, left, top, width, Emu(50800), fill=accent)
    # card body
    card = add_rect(slide, left, top + Emu(50800), width, height - Emu(50800),
                    fill=BG_CARD, line=BORDER, line_w=1.5)
    return card


# =============================================================================
# Slide 1 — Title / Hook
# =============================================================================
s = prs.slides.add_slide(BLANK)
slide_background(s, BG_DEEP)
l_corners(s)

# decorative bg panel
add_rect(s, Inches(0.8), Inches(1.2), Inches(11.7), Inches(5.1),
         fill=BG_PANEL, line=BORDER_BOLT, line_w=2)
# top accent bar
add_rect(s, Inches(0.8), Inches(1.2), Inches(11.7), Emu(63500), fill=NEON_PURPLE)

add_text(s, Inches(1.2), Inches(1.5), Inches(10), Inches(0.4),
         "▸ VIBEX://HACKATHON_PITCH", font=PIXEL_FONT, size=12,
         color=NEON_GREEN, bold=True)

# Big title
add_text(s, Inches(1.2), Inches(2.1), Inches(11), Inches(1.5),
         "VIBEX", font=PIXEL_FONT, size=96, color=NEON_PURPLE, bold=True)

# Tagline
add_text(s, Inches(1.2), Inches(3.75), Inches(11), Inches(0.7),
         "TURN YOUR AI PROJECT INTO A", font=PIXEL_FONT, size=18,
         color=TEXT, bold=True)
add_text(s, Inches(1.2), Inches(4.2), Inches(11), Inches(0.7),
         "VIRAL, EVOLVING COLLECTIBLE.", font=PIXEL_FONT, size=18,
         color=NEON_YELLOW, bold=True)

# Sub
add_text(s, Inches(1.2), Inches(5.1), Inches(11), Inches(0.5),
         "The growth layer for AI-native creators.",
         font=BODY_FONT, size=20, color=MUTED)

# bottom strip with url
add_rect(s, Inches(1.2), Inches(5.85), Inches(4.6), Inches(0.4),
         fill=NEON_PURPLE)
add_text(s, Inches(1.25), Inches(5.87), Inches(4.5), Inches(0.4),
         "  ▶  WWW.VIBEXFORGE.COM", font=PIXEL_FONT, size=11,
         color=BG_DEEP, bold=True)

add_text(s, Inches(6.0), Inches(5.9), Inches(7), Inches(0.4),
         "— by Orallexa · built with vibe coding energy",
         font=BODY_FONT, size=14, color=DIM)

add_text(s, Inches(0.5), Inches(7.05), Inches(12), Inches(0.3),
         "PRESS [ START ] TO BEGIN", font=PIXEL_FONT, size=9,
         color=NEON_GREEN, bold=True, align=PP_ALIGN.CENTER)

# =============================================================================
# Slide 2 — The Problem
# =============================================================================
s = prs.slides.add_slide(BLANK)
slide_background(s)
chrome(s, page=2)
section_title(s, "▸ CHAPTER 01 — THE PROBLEM", "ATTENTION IS THE NEW BOTTLENECK.",
              accent=NEON_ORANGE)

add_text(s, Inches(0.5), Inches(2.85), Inches(12.3), Inches(0.7),
         "Building AI products has never been easier.",
         font=BODY_FONT, size=22, color=TEXT)
add_text(s, Inches(0.5), Inches(3.3), Inches(12.3), Inches(0.7),
         "Getting attention has never been harder.",
         font=BODY_FONT, size=22, color=NEON_ORANGE, bold=True)

# 4 pain cards
pains = [
    ("NO DISTRIBUTION", "Build, ship, silence.\nNo channel after launch day."),
    ("48-HOUR HALF-LIFE", "Projects vanish from feeds\nin under two days."),
    ("NO GROWTH LOOP", "One-time spike, then zero.\nNo reason to come back."),
    ("ONE-SHOT LAUNCH", "Product Hunt happens once.\nGrowth needs a flywheel."),
]
card_w = Inches(2.95)
gap = Inches(0.2)
start_x = Inches(0.5)
top = Inches(4.35)
h = Inches(2.15)
colors = [NEON_ORANGE, NEON_PURPLE, NEON_CYAN, NEON_YELLOW]
for i, (title, body) in enumerate(pains):
    x = start_x + (card_w + gap) * i
    glass_card(s, x, top, card_w, h, accent=colors[i])
    add_text(s, x + Inches(0.2), top + Inches(0.25), card_w - Inches(0.4), Inches(0.4),
             title, font=PIXEL_FONT, size=13, color=colors[i], bold=True)
    add_text(s, x + Inches(0.2), top + Inches(0.8), card_w - Inches(0.4), Inches(1.2),
             body, font=BODY_FONT, size=15, color=TEXT)

# =============================================================================
# Slide 3 — The Insight
# =============================================================================
s = prs.slides.add_slide(BLANK)
slide_background(s)
chrome(s, page=3)
section_title(s, "▸ CHAPTER 02 — THE INSIGHT",
              "WHAT IF YOUR PROJECT LEVELED UP?", accent=NEON_GREEN)

add_text(s, Inches(0.5), Inches(2.9), Inches(12.3), Inches(0.6),
         "Games solved retention decades ago.",
         font=BODY_FONT, size=22, color=TEXT)
add_text(s, Inches(0.5), Inches(3.4), Inches(12.3), Inches(0.6),
         "Launches didn't.",
         font=BODY_FONT, size=22, color=NEON_YELLOW, bold=True)

# Big pull quote
add_rect(s, Inches(0.5), Inches(4.2), Inches(12.3), Inches(2.4),
         fill=BG_PANEL, line=NEON_PURPLE, line_w=2)
add_rect(s, Inches(0.5), Inches(4.2), Inches(12.3), Emu(50800), fill=NEON_PURPLE)

add_text(s, Inches(0.9), Inches(4.45), Inches(11.5), Inches(0.5),
         "▸ THE VIBEX INSIGHT", font=PIXEL_FONT, size=11,
         color=NEON_PURPLE, bold=True)
add_text(s, Inches(0.9), Inches(4.9), Inches(11.5), Inches(1.5),
         "Every AI project should be a LIVING ASSET —\n"
         "one that evolves, gets rarer, and pulls its own traffic.",
         font=BODY_FONT, size=24, color=TEXT, bold=True)
add_text(s, Inches(0.9), Inches(6.15), Inches(11.5), Inches(0.5),
         "Pokemon + Product Hunt + Loot Boxes, for AI creators.",
         font=BODY_FONT, size=17, color=NEON_CYAN)

# =============================================================================
# Slide 4 — Introducing VibeX
# =============================================================================
s = prs.slides.add_slide(BLANK)
slide_background(s)
chrome(s, page=4)
section_title(s, "▸ CHAPTER 03 — THE PRODUCT", "INTRODUCING VIBEX.",
              accent=NEON_PURPLE)

add_text(s, Inches(0.5), Inches(2.85), Inches(12.3), Inches(0.7),
         "A gamified growth platform for AI creators.",
         font=BODY_FONT, size=22, color=TEXT)
add_text(s, Inches(0.5), Inches(3.35), Inches(12.3), Inches(0.7),
         "We turn AI projects into shareable, evolving Hero Cards.",
         font=BODY_FONT, size=18, color=MUTED)

# Screenshot
try:
    pic = s.shapes.add_picture(
        str(SCREENS / "01-home.png"),
        Inches(0.5), Inches(4.2), width=Inches(6.1),
    )
    # Pixel frame around screenshot
    add_rect(s, Inches(0.5), Inches(4.2), Inches(6.1),
             Inches(pic.height / Inches(1)), line=NEON_PURPLE, line_w=2.5)
except Exception:
    pass

# Right side — what it is bullets
glass_card(s, Inches(7.0), Inches(4.2), Inches(5.8), Inches(2.55),
           accent=NEON_CYAN)
add_text(s, Inches(7.2), Inches(4.45), Inches(5.5), Inches(0.4),
         "▸ WHAT IT IS", font=PIXEL_FONT, size=12, color=NEON_CYAN, bold=True)

bullets = [
    ("◆", "Pixel Hero Cards for every AI project"),
    ("◆", "Stats evolve from real engagement"),
    ("◆", "QR-coded, built for social export"),
    ("◆", "Creator + VC dashboards on top"),
]
for i, (b, t) in enumerate(bullets):
    y = Inches(4.85 + i * 0.42)
    add_text(s, Inches(7.2), y, Inches(0.3), Inches(0.35),
             b, font=PIXEL_FONT, size=14, color=NEON_YELLOW, bold=True)
    add_text(s, Inches(7.55), y, Inches(5.2), Inches(0.35),
             t, font=BODY_FONT, size=16, color=TEXT)

# =============================================================================
# Slide 5 — How It Works (4 steps)
# =============================================================================
s = prs.slides.add_slide(BLANK)
slide_background(s)
chrome(s, page=5)
section_title(s, "▸ CHAPTER 04 — HOW IT WORKS", "FOUR STEPS. ZERO FRICTION.",
              accent=NEON_GREEN)

steps = [
    ("01", "CREATE", "Drop in your AI project —\nurl, readme, or idea.", NEON_PURPLE),
    ("02", "GENERATE", "Instant Hero Card\nwith HP / MP / EXP.", NEON_CYAN),
    ("03", "EVOLVE", "Level up on real\nplays, likes, shares.", NEON_GREEN),
    ("04", "SHARE", "QR-coded export\nto any social feed.", NEON_YELLOW),
]
card_w = Inches(2.95)
gap = Inches(0.2)
start_x = Inches(0.5)
top = Inches(3.2)
h = Inches(3.4)
for i, (num, title, body, c) in enumerate(steps):
    x = start_x + (card_w + gap) * i
    # card body
    glass_card(s, x, top, card_w, h, accent=c)
    # big step number
    add_text(s, x + Inches(0.2), top + Inches(0.3), card_w - Inches(0.4), Inches(1.0),
             num, font=PIXEL_FONT, size=48, color=c, bold=True)
    # title
    add_text(s, x + Inches(0.2), top + Inches(1.55), card_w - Inches(0.4), Inches(0.5),
             title, font=PIXEL_FONT, size=18, color=TEXT, bold=True)
    # divider
    add_rect(s, x + Inches(0.2), top + Inches(2.05), Inches(0.6), Emu(25400), fill=c)
    # body
    add_text(s, x + Inches(0.2), top + Inches(2.2), card_w - Inches(0.4), Inches(1.1),
             body, font=BODY_FONT, size=15, color=MUTED)
    # arrow between
    if i < 3:
        add_text(s, x + card_w - Inches(0.05), top + Inches(1.35), Inches(0.35), Inches(0.5),
                 "▶", font=PIXEL_FONT, size=22, color=c, bold=True)

# =============================================================================
# Slide 6 — Growth Loop
# =============================================================================
s = prs.slides.add_slide(BLANK)
slide_background(s)
chrome(s, page=6)
section_title(s, "▸ CHAPTER 05 — THE MOAT", "BUILT-IN DISTRIBUTION ENGINE.",
              accent=NEON_YELLOW)

# Loop nodes arranged horizontally
nodes = ["CREATE", "CARD", "SHARE", "TRAFFIC", "LEVEL UP"]
nc = [NEON_PURPLE, NEON_CYAN, NEON_PINK, NEON_YELLOW, NEON_GREEN]
node_w = Inches(2.1)
node_h = Inches(1.0)
y = Inches(3.7)
total_w = len(nodes) * node_w + (len(nodes) - 1) * Inches(0.35)
start = (SLIDE_W - total_w) / 2

for i, label in enumerate(nodes):
    x = start + i * (node_w + Inches(0.35))
    # outer neon border box
    add_rect(s, x, y, node_w, node_h, fill=BG_CARD, line=nc[i], line_w=3)
    add_rect(s, x, y, node_w, Emu(38100), fill=nc[i])
    add_text(s, x, y + Inches(0.3), node_w, Inches(0.5),
             label, font=PIXEL_FONT, size=16, color=TEXT, bold=True,
             align=PP_ALIGN.CENTER)
    if i < len(nodes) - 1:
        add_text(s, x + node_w, y + Inches(0.2), Inches(0.35), Inches(0.6),
                 "▶", font=PIXEL_FONT, size=22, color=NEON_GREEN, bold=True,
                 align=PP_ALIGN.CENTER)

# return arrow (loop)
add_text(s, Inches(0.5), Inches(4.95), Inches(12.3), Inches(0.5),
         "↑ ————————————— SHARE AGAIN ————————————— ↓",
         font=PIXEL_FONT, size=12, color=NEON_YELLOW, bold=True,
         align=PP_ALIGN.CENTER)

# caption
add_text(s, Inches(0.5), Inches(5.75), Inches(12.3), Inches(0.6),
         "Every card is an ad for itself.",
         font=BODY_FONT, size=22, color=TEXT, bold=True,
         align=PP_ALIGN.CENTER)
add_text(s, Inches(0.5), Inches(6.25), Inches(12.3), Inches(0.5),
         "Product Hunt gives you 48 hours. VibeX gives you a flywheel.",
         font=BODY_FONT, size=17, color=MUTED, align=PP_ALIGN.CENTER)

# =============================================================================
# Slide 7 — Product Tour (screenshots)
# =============================================================================
s = prs.slides.add_slide(BLANK)
slide_background(s)
chrome(s, page=7)
section_title(s, "▸ CHAPTER 06 — SEE IT LIVE", "THE APP, RIGHT NOW.",
              accent=NEON_CYAN)

shots = [
    ("03-discover.png", "DISCOVER", NEON_PURPLE),
    ("05-feed.png", "FEED", NEON_CYAN),
    ("08-vc.png", "VC DASHBOARD", NEON_YELLOW),
]
sw = Inches(4.0)
gap = Inches(0.25)
sx0 = Inches(0.5)
sy = Inches(3.1)
for i, (fname, label, c) in enumerate(shots):
    sx = sx0 + (sw + gap) * i
    try:
        pic = s.shapes.add_picture(
            str(SCREENS / fname), sx, sy, width=sw,
        )
        ph = Inches(pic.height / Inches(1))
    except Exception:
        ph = Inches(2.5)
    # label bar under
    lby = sy + ph + Inches(0.15)
    add_rect(s, sx, lby, sw, Inches(0.4), fill=BG_CARD, line=c, line_w=2)
    add_rect(s, sx, lby, Inches(0.15), Inches(0.4), fill=c)
    add_text(s, sx + Inches(0.3), lby + Inches(0.07), sw - Inches(0.4), Inches(0.3),
             label, font=PIXEL_FONT, size=12, color=TEXT, bold=True)

add_text(s, Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.45),
         "48 tables · 43 API routes · 105+ components · shipping weekly.",
         font=BODY_FONT, size=16, color=MUTED, align=PP_ALIGN.CENTER)

# =============================================================================
# Slide 8 — Evolution System
# =============================================================================
s = prs.slides.add_slide(BLANK)
slide_background(s)
chrome(s, page=8)
section_title(s, "▸ CHAPTER 07 — THE GAME LAYER",
              "6-STAGE EVOLUTION SYSTEM.", accent=NEON_PURPLE)

rows = [
    ("SEED",     "COMMON",    "Just created",                         NEON_GREEN),
    ("ACTIVE",   "UNCOMMON",  "50+ plays or 40+ score",               NEON_CYAN),
    ("GROWING",  "RARE",      "500+ plays · 50+ upvotes · 1K+ views", NEON_PURPLE),
    ("BREAKOUT", "EPIC",      "1K+ plays · 100+ upvotes · 70+ orig.", NEON_ORANGE),
    ("LEGEND",   "LEGENDARY", "5K+ plays · 500+ upvotes · 85+ score", NEON_YELLOW),
    ("MYTH",     "MYTHIC",    "10K+ plays · 1K+ upvotes · VC interest", NEON_PINK),
]

top = Inches(2.9)
row_h = Inches(0.6)
x0 = Inches(0.5)
tw = Inches(12.3)

# header
add_rect(s, x0, top, tw, row_h, fill=BG_PANEL, line=BORDER_BOLT, line_w=1.5)
add_text(s, x0 + Inches(0.25), top + Inches(0.15), Inches(2.5), Inches(0.3),
         "STAGE", font=PIXEL_FONT, size=11, color=NEON_GREEN, bold=True)
add_text(s, x0 + Inches(3), top + Inches(0.15), Inches(2.5), Inches(0.3),
         "RARITY", font=PIXEL_FONT, size=11, color=NEON_GREEN, bold=True)
add_text(s, x0 + Inches(6), top + Inches(0.15), Inches(6), Inches(0.3),
         "TRIGGER", font=PIXEL_FONT, size=11, color=NEON_GREEN, bold=True)

for i, (stage, rarity, trig, c) in enumerate(rows):
    y = top + row_h + row_h * i
    add_rect(s, x0, y, tw, row_h, fill=BG_CARD, line=BORDER, line_w=1)
    # color tick
    add_rect(s, x0, y, Emu(50800), row_h, fill=c)
    add_text(s, x0 + Inches(0.25), y + Inches(0.13), Inches(2.7), Inches(0.35),
             stage, font=PIXEL_FONT, size=13, color=c, bold=True)
    add_text(s, x0 + Inches(3), y + Inches(0.13), Inches(2.8), Inches(0.35),
             rarity, font=PIXEL_FONT, size=12, color=TEXT, bold=True)
    add_text(s, x0 + Inches(6), y + Inches(0.15), Inches(6), Inches(0.35),
             trig, font=BODY_FONT, size=15, color=MUTED)

# =============================================================================
# Slide 9 — Why VibeX (comparison)
# =============================================================================
s = prs.slides.add_slide(BLANK)
slide_background(s)
chrome(s, page=9)
section_title(s, "▸ CHAPTER 08 — WHY VIBEX", "WE'RE NOT ANOTHER LAUNCH LIST.",
              accent=NEON_ORANGE)

rivals = [
    ("PRODUCT HUNT", "One-time launch",       NEON_ORANGE, False),
    ("GITHUB",       "No discovery",          NEON_ORANGE, False),
    ("X / TWITTER",  "High noise",            NEON_ORANGE, False),
    ("VIBEX",        "Continuous growth loop", NEON_GREEN, True),
]

card_w = Inches(2.95)
gap = Inches(0.2)
x0 = Inches(0.5)
y = Inches(3.35)
h = Inches(3.0)

for i, (name, desc, c, win) in enumerate(rivals):
    x = x0 + (card_w + gap) * i
    fill = BG_PANEL if win else BG_CARD
    line_w = 3 if win else 1.5
    add_rect(s, x, y, card_w, h, fill=fill, line=c, line_w=line_w)
    # top bar
    add_rect(s, x, y, card_w, Emu(50800), fill=c)
    # badge
    icon = "✓" if win else "✗"
    add_text(s, x + Inches(0.2), y + Inches(0.3), card_w - Inches(0.4), Inches(0.8),
             icon, font=PIXEL_FONT, size=48, color=c, bold=True,
             align=PP_ALIGN.CENTER)
    # name
    add_text(s, x + Inches(0.2), y + Inches(1.4), card_w - Inches(0.4), Inches(0.5),
             name, font=PIXEL_FONT, size=14, color=TEXT, bold=True,
             align=PP_ALIGN.CENTER)
    # divider
    add_rect(s, x + card_w / 2 - Inches(0.3), y + Inches(1.95), Inches(0.6), Emu(25400),
             fill=c)
    # desc
    add_text(s, x + Inches(0.2), y + Inches(2.1), card_w - Inches(0.4), Inches(0.8),
             desc, font=BODY_FONT, size=15, color=MUTED,
             align=PP_ALIGN.CENTER)

add_text(s, Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.4),
         "Others give you a moment. VibeX gives you a mechanic.",
         font=BODY_FONT, size=16, color=NEON_YELLOW, bold=True,
         align=PP_ALIGN.CENTER)

# =============================================================================
# Slide 10 — Tech Stack
# =============================================================================
s = prs.slides.add_slide(BLANK)
slide_background(s)
chrome(s, page=10)
section_title(s, "▸ CHAPTER 09 — UNDER THE HOOD", "BUILT ON MODERN RAILS.",
              accent=NEON_CYAN)

stack = [
    ("FRAMEWORK",  "Next.js 16 · Turbopack · App Router", NEON_PURPLE),
    ("LANGUAGE",   "TypeScript 5 (strict)",               NEON_CYAN),
    ("UI",         "React 19 · Tailwind 4 · Framer",      NEON_PINK),
    ("DESIGN",     "NES.css + RPGUI · 16-bit pixel",      NEON_YELLOW),
    ("DATABASE",   "Supabase · Postgres · RLS · Realtime", NEON_GREEN),
    ("AUTH",       "Supabase Auth · GitHub + Google",     NEON_CYAN),
    ("AI",         "Claude API · launch · reviews · growth", NEON_PURPLE),
    ("MONITORING", "Sentry · PostHog",                    NEON_ORANGE),
]

col_w = Inches(6.1)
row_h = Inches(0.7)
x_left = Inches(0.5)
x_right = Inches(6.75)
y0 = Inches(3.0)

for i, (label, val, c) in enumerate(stack):
    col = i % 2
    row = i // 2
    x = x_left if col == 0 else x_right
    y = y0 + row_h * row + Inches(0.15) * row
    add_rect(s, x, y, col_w, row_h, fill=BG_CARD, line=BORDER, line_w=1)
    add_rect(s, x, y, Emu(50800), row_h, fill=c)
    add_text(s, x + Inches(0.25), y + Inches(0.17), Inches(2), Inches(0.35),
             label, font=PIXEL_FONT, size=11, color=c, bold=True)
    add_text(s, x + Inches(2.4), y + Inches(0.2), col_w - Inches(2.6), Inches(0.35),
             val, font=BODY_FONT, size=15, color=TEXT)

add_text(s, Inches(0.5), Inches(6.75), Inches(12.3), Inches(0.4),
         "`npm run dev` → runs in 30 seconds · zero config · mock mode baked in.",
         font=BODY_FONT, size=15, color=NEON_GREEN, align=PP_ALIGN.CENTER)

# =============================================================================
# Slide 11 — Monetization
# =============================================================================
s = prs.slides.add_slide(BLANK)
slide_background(s)
chrome(s, page=11)
section_title(s, "▸ CHAPTER 10 — BUSINESS MODEL", "HOW VIBEX MAKES MONEY.",
              accent=NEON_YELLOW)

rev = [
    ("CARD UPGRADES",    "Rare / Epic / Legendary\nvisual skins.",      NEON_PURPLE),
    ("FEATURED SLOTS",   "Paid placement on the\ndiscover page.",       NEON_CYAN),
    ("GROWTH INSIGHTS",  "AI-powered analytics\nfor serious creators.", NEON_GREEN),
    ("AGENT MARKET",     "Premium agents with\nrev-share for authors.", NEON_PINK),
]

card_w = Inches(2.95)
gap = Inches(0.2)
x0 = Inches(0.5)
top = Inches(3.25)
h = Inches(3.1)
for i, (title, body, c) in enumerate(rev):
    x = x0 + (card_w + gap) * i
    glass_card(s, x, top, card_w, h, accent=c)
    add_text(s, x + Inches(0.2), top + Inches(0.35), card_w - Inches(0.4), Inches(0.5),
             f"0{i+1}", font=PIXEL_FONT, size=26, color=c, bold=True)
    add_text(s, x + Inches(0.2), top + Inches(1.05), card_w - Inches(0.4), Inches(0.5),
             title, font=PIXEL_FONT, size=13, color=TEXT, bold=True)
    add_rect(s, x + Inches(0.2), top + Inches(1.55), Inches(0.6), Emu(25400),
             fill=c)
    add_text(s, x + Inches(0.2), top + Inches(1.7), card_w - Inches(0.4), Inches(1.3),
             body, font=BODY_FONT, size=15, color=MUTED)

add_text(s, Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.4),
         "Growth loop first. Monetization rides the loop, doesn't gate it.",
         font=BODY_FONT, size=15, color=NEON_YELLOW, bold=True,
         align=PP_ALIGN.CENTER)

# =============================================================================
# Slide 12 — Traction / Status
# =============================================================================
s = prs.slides.add_slide(BLANK)
slide_background(s)
chrome(s, page=12)
section_title(s, "▸ CHAPTER 11 — WHERE WE ARE", "SHIPPED. LIVE. ITERATING.",
              accent=NEON_GREEN)

stats = [
    ("48",   "DB TABLES",      NEON_PURPLE),
    ("43",   "API ROUTES",     NEON_CYAN),
    ("105+", "COMPONENTS",     NEON_YELLOW),
    ("708",  "i18n KEYS",      NEON_GREEN),
]
sw = Inches(2.95)
gap = Inches(0.2)
x0 = Inches(0.5)
y = Inches(3.0)
h = Inches(1.9)
for i, (num, label, c) in enumerate(stats):
    x = x0 + (sw + gap) * i
    add_rect(s, x, y, sw, h, fill=BG_CARD, line=c, line_w=2)
    add_rect(s, x, y, sw, Emu(50800), fill=c)
    add_text(s, x, y + Inches(0.3), sw, Inches(1.0),
             num, font=PIXEL_FONT, size=48, color=c, bold=True,
             align=PP_ALIGN.CENTER)
    add_text(s, x, y + Inches(1.35), sw, Inches(0.4),
             label, font=PIXEL_FONT, size=12, color=TEXT, bold=True,
             align=PP_ALIGN.CENTER)

# Status panel
add_rect(s, Inches(0.5), Inches(5.2), Inches(12.3), Inches(1.5),
         fill=BG_PANEL, line=NEON_GREEN, line_w=2)
add_rect(s, Inches(0.5), Inches(5.2), Inches(12.3), Emu(50800), fill=NEON_GREEN)
add_text(s, Inches(0.7), Inches(5.4), Inches(6), Inches(0.4),
         "▸ STATUS: EARLY MVP · LIVE · OPEN BETA",
         font=PIXEL_FONT, size=12, color=NEON_GREEN, bold=True)
add_text(s, Inches(0.7), Inches(5.8), Inches(12), Inches(0.45),
         "Real Supabase + Claude API · i18n EN/ZH ·",
         font=BODY_FONT, size=17, color=TEXT)
add_text(s, Inches(0.7), Inches(6.15), Inches(12), Inches(0.45),
         "Source-available on GitHub · try it at vibexforge.com",
         font=BODY_FONT, size=17, color=NEON_YELLOW)

# =============================================================================
# Slide 13 — Vision
# =============================================================================
s = prs.slides.add_slide(BLANK)
slide_background(s)
chrome(s, page=13)
section_title(s, "▸ CHAPTER 12 — THE VISION", "", accent=NEON_PURPLE)

# Big statement
add_text(s, Inches(0.5), Inches(2.2), Inches(12.3), Inches(1.0),
         "BECOME THE GROWTH LAYER",
         font=PIXEL_FONT, size=34, color=TEXT, bold=True,
         align=PP_ALIGN.CENTER)
add_text(s, Inches(0.5), Inches(3.0), Inches(12.3), Inches(1.0),
         "FOR AI CREATORS.",
         font=PIXEL_FONT, size=34, color=NEON_PURPLE, bold=True,
         align=PP_ALIGN.CENTER)

# Divider
add_rect(s, Inches(5.8), Inches(4.0), Inches(1.7), Emu(38100),
         fill=NEON_PURPLE)

add_text(s, Inches(0.5), Inches(4.3), Inches(12.3), Inches(0.6),
         "Every project is not just launched —",
         font=BODY_FONT, size=24, color=TEXT,
         align=PP_ALIGN.CENTER)
add_text(s, Inches(0.5), Inches(4.8), Inches(12.3), Inches(0.6),
         "it evolves.",
         font=BODY_FONT, size=28, color=NEON_YELLOW, bold=True,
         align=PP_ALIGN.CENTER)

# Pixel footer decoration
add_text(s, Inches(0.5), Inches(5.95), Inches(12.3), Inches(0.5),
         "◆  ◆  ◆",
         font=PIXEL_FONT, size=20, color=NEON_PURPLE, bold=True,
         align=PP_ALIGN.CENTER)

add_text(s, Inches(0.5), Inches(6.45), Inches(12.3), Inches(0.5),
         "AI makes building easy. VibeX makes growing inevitable.",
         font=BODY_FONT, size=18, color=MUTED, align=PP_ALIGN.CENTER)

# =============================================================================
# Slide 14 — Thank You / CTA
# =============================================================================
s = prs.slides.add_slide(BLANK)
slide_background(s)
l_corners(s)

# Full-bleed dark panel
add_rect(s, Inches(0.8), Inches(0.9), Inches(11.7), Inches(5.7),
         fill=BG_PANEL, line=BORDER_BOLT, line_w=2)
add_rect(s, Inches(0.8), Inches(0.9), Inches(11.7), Emu(63500), fill=NEON_GREEN)

add_text(s, Inches(1.2), Inches(1.15), Inches(11), Inches(0.4),
         "▸ VIBEX://END_OF_PITCH", font=PIXEL_FONT, size=12,
         color=NEON_GREEN, bold=True)

# THANK YOU huge
add_text(s, Inches(1.2), Inches(1.7), Inches(11), Inches(1.5),
         "THANK YOU.", font=PIXEL_FONT, size=80, color=TEXT, bold=True,
         align=PP_ALIGN.CENTER)

add_text(s, Inches(1.2), Inches(3.3), Inches(11), Inches(0.5),
         "NOW GO EVOLVE SOMETHING.",
         font=PIXEL_FONT, size=18, color=NEON_YELLOW, bold=True,
         align=PP_ALIGN.CENTER)

# CTA box
cta_w = Inches(6.0)
cta_x = (SLIDE_W - cta_w) / 2
add_rect(s, cta_x, Inches(4.2), cta_w, Inches(0.7),
         fill=NEON_PURPLE, line=TEXT, line_w=2)
add_text(s, cta_x, Inches(4.35), cta_w, Inches(0.4),
         "▶  TRY IT · WWW.VIBEXFORGE.COM",
         font=PIXEL_FONT, size=15, color=BG_DEEP, bold=True,
         align=PP_ALIGN.CENTER)

add_rect(s, cta_x, Inches(5.05), cta_w, Inches(0.7),
         fill=BG_CARD, line=NEON_GREEN, line_w=2)
add_text(s, cta_x, Inches(5.2), cta_w, Inches(0.4),
         "★  STAR ON GITHUB · ALEX-JB/VIBEX",
         font=PIXEL_FONT, size=15, color=NEON_GREEN, bold=True,
         align=PP_ALIGN.CENTER)

add_text(s, Inches(1.2), Inches(6.1), Inches(11), Inches(0.4),
         "Built with vibe coding energy — by Orallexa",
         font=BODY_FONT, size=16, color=MUTED, align=PP_ALIGN.CENTER)

add_text(s, Inches(0.5), Inches(7.05), Inches(12), Inches(0.3),
         "[ SELECT ]     [ START ]", font=PIXEL_FONT, size=9,
         color=DIM, bold=True, align=PP_ALIGN.CENTER)

# -----------------------------------------------------------------------------
prs.save(str(OUT))
print(f"Saved → {OUT}")
print(f"Slides: {len(prs.slides)}")
